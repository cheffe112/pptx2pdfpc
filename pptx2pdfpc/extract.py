import json
import pathlib
from typing import Tuple, List, Any

from pptx import Presentation
from pptx.exc import PackageNotFoundError

# logic for extracting speaker notes and text boxes by user fusion on Stack Overflow, see:
# https://stackoverflow.com/questions/63659972/extract-presenter-notes-from-pptx-file-powerpoint
# Many thanks <3
from pptx2pdfpc.errors import UsageError


def speaker_notes(input_pptx: pathlib.Path) -> List[Tuple[int, str]]:
    """Extract all speaker notes from a pptx."""
    try:
        prs = Presentation(str(input_pptx))
    except PackageNotFoundError as e:
        raise UsageError(f"There was an error while reading in the pptx:\n{e}")

    extracted = []
    for page, slide in enumerate(prs.slides, start=1):
        text = slide.notes_slide.notes_text_frame.text
        extracted.append((page, text))
    return extracted


def text_boxes(input_pptx: pathlib.Path) -> List[Tuple[int, List[Any]]]:
    """Extract all text boxes from all slides."""
    prs = Presentation(str(input_pptx))
    extracted = []
    for page, slide in enumerate(prs.slides, start=1):
        temp = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                temp.append(shape.text)
        extracted.append((page, temp))
    return extracted


def generate_pdfpc(extracted_notes: List[Tuple[int, str]], options: List[Tuple[str, Any]], output_path: pathlib.Path,
                   font_size: str = None):
    """Generate a config file for pdfpc and writes it to output_path.
    The file ending is .pdfpc and can be used together with a PDF
    version of the pptx presentation.
    From the pdfpc man page: "When pdfpc is invoked with a PDF file, it automatically
    checks for and loads the associated .pdfpc file, if it exists."

    Requirement: Both files need to have the same name, one ending in .pdf, the other in .pdfpc.
    For pdfpc, see its man page or https://man.archlinux.org/man/community/pdfpc/pdfpc.1.en
    """
    options_to_write = _create_options(options)

    DELIMITER = "###"
    with output_path.open("w") as fo:
        # fo.writelines(options_to_write)
        pages = []
        idx = 0
        for slide in extracted_notes:
            page = {'idx': idx, "label": str(idx + 1), "overlay": 0}
            idx = idx + 1
            page_number = slide[0]
            note_text = slide[1]
            if len(note_text) > 0:
                print(note_text)
                page["note"] = note_text
            pages.append(page)

        dataset = {"pdfpcFormat": 1, "disableMarkdown": True, "noteFontSize": 20, "pages": pages}
        json.dump(dataset, fo, indent=6)


def generate_output_path(input_path: pathlib.Path) -> pathlib.Path:
    """Generate the path to the notes file ending in .pdfpc. It must have
    the same name as the presentation pdf and be in the same directory
    to be visible during the presentation."""
    filename = input_path.stem
    path = input_path.parent
    return path / pathlib.Path(filename + ".pdfpc")


def _create_options(options: List[Tuple[str, Any]]) -> List[str]:
    """Parse cli options and convert them to .pdfpc options

    A list of available options .pdfpc options can be found here:
    https://github.com/pdfpc/pdfpc/pull/556#issuecomment-739565800
    duration 			 -d, --duration=N           Duration in minutes of the presentation used for timer  display.
    start_time 		 -t, --start-time=T         Start time of the presentation to be used as a countdown.  (Format: HH:MM (24h))
    end_time 			 -e, --end-time=T           End time of the presentation. (Format: HH:MM (24h))
    last_minutes 	5 	 -l, --last-minutes=N       Time in minutes, from which on the timer changes its color. (Default: 5 minutes)
    notes_position 	 -p  --notes-position       Position  of notes on the PDF page. Position can be either left, right, top or  bottom.
    default_transition -a --default-transition
    disable_markdown 	 -m --disable-markdown      FALSE
    end_slide 		 -s --end-slide             Int.
    last_saved_slide 	 -i --last-saved-slide
    font_size 		 -f --font-size             Font size in pt. Int.
    notes 			 -n --notes
    forced_overlays 	 -v --forced-overlays
    the options in pptx2pdfpc follow the command-line options in pdfpc if available
    """

    options_to_write = []
    for option in options:
        if option[1]:
            options_to_write.append(f"[{option[0]}]\n{option[1]}\n")

    return options_to_write
